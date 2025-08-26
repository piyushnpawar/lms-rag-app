from langchain_community.document_loaders import PyMuPDFLoader
from langchain.docstore.document import Document as LangchainDocument
from langchain_text_splitters import RecursiveCharacterTextSplitter
from langchain_community.embeddings.fastembed import FastEmbedEmbeddings
from langchain_qdrant import QdrantVectorStore
from qdrant_client import QdrantClient
from qdrant_client.http.models import Distance, VectorParams
import hashlib, os, uuid, logging, requests, magic
import lms_handling
from docx import Document
from pptx import Presentation
from dotenv import load_dotenv
load_dotenv()

QDRANT_COLLECTION_NAME = "LMS"
QDRANT_URL = os.environ.get("QDRANT_URL")
EMBEDDING_MODEL = FastEmbedEmbeddings(model_name="BAAI/bge-small-en-v1.5")

try:
    client = QdrantClient(
        url=QDRANT_URL
    )
    QDRANT_INSTANCE = QdrantVectorStore(
        client=client,
        collection_name=QDRANT_COLLECTION_NAME,
        embedding=EMBEDDING_MODEL,
    )
except:
    client = QdrantClient(
        url=QDRANT_URL
    )
    client.create_collection(
        collection_name=QDRANT_COLLECTION_NAME,
        vectors_config=VectorParams(size=384, distance=Distance.COSINE),
    )
    QDRANT_INSTANCE = QdrantVectorStore(
        client=client,
        collection_name=QDRANT_COLLECTION_NAME,
        embedding=EMBEDDING_MODEL,
    )

SUPPORTED_MIME_TYPES = {
    'application/pdf',
    'application/vnd.openxmlformats-officedocument.wordprocessingml.document',
    'application/vnd.openxmlformats-officedocument.presentationml.presentation'
}


async def ingestData(subject: str, file_name:str, file_source:str) -> str:
    session = requests.Session()
    requests.utils.add_dict_to_cookiejar(session.cookies, lms_handling.SESSION_COOKIES)
    try:
        logging.info(f"Getting file: {file_name} from {file_source}")
        response = session.get(file_source)
        response.raise_for_status()
        document_content = response.content
        document_hash = calculate_pdf_hash(document_content)
        logging.info(f"Calculated PDF hash: {document_hash}")
    except requests.RequestException as e:
        logging.error(f"Failed to download file: {e}")
        return "Source invalid"

    logging.info("Checking if the content already exists")
    if await check_for_existing_document(document_hash):
        return "Already exists in database"

    check, docs = extract_document(document_content, document_hash, subject, file_name, file_source)
    if check:
        text_splitter = RecursiveCharacterTextSplitter(chunk_size=1000, chunk_overlap=300)
        split_docs = text_splitter.split_documents(documents=docs)
        logging.info("Document chunking complete")
    else:
        logging.info("Document extraction failed")
        return docs

    try:
        await QDRANT_INSTANCE.aadd_documents(split_docs)
        logging.info("Chunks loaded in DB")
        return "File loaded successfully"
    except Exception as e:
        logging.info(f"Chunk loading failed: {e}")
        return "An error occured"

def calculate_pdf_hash(pdf_content: bytes) -> str:
    return hashlib.sha256(pdf_content).hexdigest()

def extract_document(document_content, document_hash, subject, file_name, file_source):
    logging.info("Attempting to load document")
    mime_type = magic.from_buffer(document_content, mime=True)
    temp_file_path=""
    if mime_type not in SUPPORTED_MIME_TYPES:
        logging.warning(f"Unsupported file type detected: {mime_type}. Skipping ingestion.")
        return False,"Unsupported file type"

    if mime_type == 'application/pdf':
        file_extension = ".pdf"
        temp_file_path = create_temp_file(file_extension,document_content)
        loader = PyMuPDFLoader(temp_file_path)       
        docs = loader.load()  
        for doc in docs:
            doc.metadata["content_hash"] = document_hash
            doc.metadata["subject"] = subject
            doc.metadata["file_name"] = file_name
            doc.metadata["file_link"] = file_source
    elif mime_type == 'application/vnd.openxmlformats-officedocument.wordprocessingml.document':
        file_extension = ".docx"
        temp_file_path = create_temp_file(file_extension,document_content)
        document = Document(temp_file_path)
        full_text = []
        for paragraph in document.paragraphs:
            full_text.append(paragraph.text)
        text_content = "/n".join(full_text)
        doc_metadata = {
            "content_hash": document_hash,
            "subject": subject,
            "file_name": file_name,
            "file_link": file_source
        }
        docs = [LangchainDocument(
        page_content=text_content,
        metadata=doc_metadata
        )]
    elif mime_type == 'application/vnd.openxmlformats-officedocument.presentationml.presentation':
        file_extension = ".pptx"
        temp_file_path = create_temp_file(file_extension,document_content)
        prs = Presentation(temp_file_path)
        full_text = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "has_text_frame") and shape.has_text_frame:
                    full_text.append(shape.text)
        doc_metadata = {
            "content_hash": document_hash,
            "subject": subject,
            "file_name": file_name,
            "file_link": file_source
        }
        docs = [Document(page_content="\n".join(full_text), metadata=doc_metadata)]
    else:
        logging.error(f"Internal error: Unsupported MIME type {mime_type} passed the check.")
        return False,"An internal error occurred"

    os.remove(temp_file_path)
    
    if not docs:
        logging.error("Document loading failed")
        return False, "An error occured"
    return True, docs

def create_temp_file(file_extension, document_content):
    temp_file_path = f"temp_{uuid.uuid4()}{file_extension}"
    with open(temp_file_path, "wb") as f:
        f.write(document_content)
    return temp_file_path

async def check_for_existing_document(document_hash: str) -> bool:
    try:
        qdrant_filter = {
            "must": [
                {
                    "key": "metadata.content_hash",
                    "match": {
                        "value": document_hash
                    }
                }
            ]
        }
        search_results = await QDRANT_INSTANCE.asimilarity_search("dummy query", k=1, filter=qdrant_filter)
        if search_results:
            logging.info("Document with the same source and content hash already exists. Skipping ingestion")
            return True
        else:
            logging.info("Document not found. Proceeding with ingestion")
            return False

    except Exception as e:
        logging.info(f"Collection '{QDRANT_COLLECTION_NAME}' does not exist yet. Assuming no duplicates: {e}")
        return False